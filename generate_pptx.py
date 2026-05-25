#!/usr/bin/env python3
"""
Generate PPTX from rotating-narrator JSON config.

Layout rules:
- Slide: 19.05 cm × 33.87 cm (720×1280 portrait)
- Each new line appears at a fixed bottom anchor (18.24 cm from top).
- Previous lines shift upward to make room; their display size is scaled
  proportionally to the new line's font size (camera-zoom simulation).
- Font size of each new line auto-fits text to the visible slide width.
- Completed phases rotate ±90° and are packed off-slide to the left/right.
- If there are more than 3 old phases, the oldest is discarded.
"""

import json
import re
import os
import sys
from io import BytesIO

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image, ImageFilter

# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------
SLIDE_W_CM   = 19.05
SLIDE_H_CM   = 33.87
LEFT_CM      = 2.54          # left edge of all active text
BOTTOM_CM    = 18.24         # bottom of the newest line
BASE_PT      = 48            # max font size (shortest texts)
SCALE_BACK   = 0.82          # scale factor per step backward (camera-zoom sim)
GAP_CM       = 0.35          # vertical gap between stacked lines
BOX_W_CM     = 26.0          # text-box width (extends off-slide; no wrap)
AVAIL_W_CM   = SLIDE_W_CM - LEFT_CM   # visible text width for auto-fit

MAX_OLD_PHASES = 3           # discard oldest phase when old count exceeds this


def cm(v): return Cm(v)
def pt(v): return Pt(v)


# ---------------------------------------------------------------------------
# Colour helper — supports both #RGB and #RRGGBB
# ---------------------------------------------------------------------------
def parse_color(hex_str):
    h = hex_str.lstrip('#')
    if len(h) == 3:
        h = h[0]*2 + h[1]*2 + h[2]*2
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


# ---------------------------------------------------------------------------
# Markup parser  {&font=X&color=#RRGGBB&size=+N}text{/}
# ---------------------------------------------------------------------------
MARKUP_RE = re.compile(r'\{&([^}]+)\}(.*?)\{/\}', re.DOTALL)


def parse_runs(text, default_font, default_color_hex, default_size_pt):
    """Return list of (segment_text, font_name, color_hex, size_pt)."""
    runs, last = [], 0
    for m in MARKUP_RE.finditer(text):
        if m.start() > last:
            runs.append((text[last:m.start()], default_font,
                         default_color_hex, default_size_pt))
        attrs = {}
        for part in m.group(1).split('&'):
            if '=' in part:
                k, v = part.split('=', 1)
                attrs[k] = v
        font  = attrs.get('font',  default_font)
        color = attrs.get('color', default_color_hex)
        size_raw = attrs.get('size')
        if size_raw is not None:
            try:
                size = default_size_pt + int(size_raw.lstrip('+'))
            except ValueError:
                size = default_size_pt
        else:
            size = default_size_pt
        runs.append((m.group(2), font, color, size))
        last = m.end()
    if last < len(text):
        runs.append((text[last:], default_font, default_color_hex, default_size_pt))
    return runs


# ---------------------------------------------------------------------------
# Font-size auto-fit
# ---------------------------------------------------------------------------
def strip_markup(text):
    """Return plain text with markup tags removed."""
    return re.sub(r'\{[^}]+\}', '', text)


def natural_font_size(text, max_pt=BASE_PT):
    """
    Compute font size so the text fits in AVAIL_W_CM.
    CJK characters are full-width (1 em); ASCII/punct are ~0.6 em.
    """
    plain = strip_markup(text)
    if not plain:
        return max_pt
    cjk   = sum(1 for c in plain if '⺀' <= c <= '鿿'
                                  or '豈' <= c <= '﫿')
    other = len(plain) - cjk
    em_cm = max_pt / 72.0 * 2.54
    total = (cjk + other * 0.6) * em_cm
    if total <= AVAIL_W_CM:
        return max_pt
    return max_pt * AVAIL_W_CM / total


# ---------------------------------------------------------------------------
# Background image (blurred)
# ---------------------------------------------------------------------------
def make_blurred_bg(image_path, blur_radius, w=720, h=1280):
    img = Image.open(image_path).convert('RGB')
    img = img.resize((w, h), Image.LANCZOS)
    img = img.filter(ImageFilter.GaussianBlur(radius=blur_radius))
    buf = BytesIO()
    img.save(buf, format='JPEG', quality=85)
    buf.seek(0)
    return buf


def add_background(slide, image_path, blur_radius):
    buf = make_blurred_bg(image_path, blur_radius)
    pic = slide.shapes.add_picture(buf, cm(0), cm(0), cm(SLIDE_W_CM), cm(SLIDE_H_CM))
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)


# ---------------------------------------------------------------------------
# Text-box helpers
# ---------------------------------------------------------------------------
def line_height_cm(font_size_pt):
    return font_size_pt / 72.0 * 2.54 * 1.45


def add_textbox(slide, text, default_font, default_color_hex, font_size_pt,
                left_cm, top_cm, width_cm, height_cm, rotation_deg=0):
    txBox = slide.shapes.add_textbox(
        cm(left_cm), cm(top_cm), cm(width_cm), cm(height_cm))

    if rotation_deg != 0:
        sp   = txBox._element
        spPr = sp.find(qn('p:spPr'))
        xfrm = spPr.find(qn('a:xfrm'))
        if xfrm is None:
            xfrm = etree.SubElement(spPr, qn('a:xfrm'))
        xfrm.set('rot', str(int(rotation_deg * 60000)))

    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]

    for seg_text, font_name, color_hex, size_pt in \
            parse_runs(text, default_font, default_color_hex, font_size_pt):
        if not seg_text:
            continue
        run = p.add_run()
        run.text = seg_text
        run.font.name = font_name
        run.font.size = Pt(size_pt)
        try:
            run.font.color.rgb = parse_color(color_hex)
        except Exception:
            pass

    txBody = tf._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        for tag in [qn('a:noAutofit'), qn('a:spAutoFit'), qn('a:normAutofit')]:
            el = bodyPr.find(tag)
            if el is not None:
                bodyPr.remove(el)
        etree.SubElement(bodyPr, qn('a:noAutofit'))

    return txBox


# ---------------------------------------------------------------------------
# Entrance animation
# ---------------------------------------------------------------------------
_ANIM_ID = [0]

def _next_id():
    _ANIM_ID[0] += 1
    return _ANIM_ID[0]


# Translate old Chinese direction words to canonical English params
_ZH_PARAM = {
    '向左': 'left', '向右': 'right', '向上': 'up', '向下': 'down',
    '横向': 'horizontal', '纵向': 'vertical',
    '向内': 'in', '向外': 'out',
    '水平向外': 'h-out', '水平向内': 'h-in',
    '垂直向外': 'v-out', '垂直向内': 'v-in',
}

# Translate old Chinese effect names to canonical English names
_ZH_NAME = {
    '百叶窗': 'blinds',     '出现': 'appear',      '盒状': 'box',
    '菱形': 'diamond',      '劈裂': 'split',        '切入': 'strips',
    '随机线条': 'random-bars', '楔入': 'wedge',      '擦除': 'wipe',
    '飞入': 'fly',          '阶梯状': 'stairs',     '轮子': 'wheel',
    '棋盘': 'checkerboard', '十字形扩展': 'plus',   '向内溶解': 'dissolve',
    '圆形扩展': 'circle',   '淡化': 'fade',         '旋转': 'spin',
    '缩放': 'zoom',         '展开': 'expand',       '翻转式由远及近': 'flip',
    '基本缩放': 'zoom-basic', '伸展': 'stretch',    '下浮': 'float-down',
    '回旋': 'swivel',       '上浮': 'float-up',     '升起': 'rise',
    '压缩': 'compress',     '中心旋转': 'center-revolve', '弹跳': 'bounce',
    '飞旋': 'pinwheel',     '挥鞭式': 'whip',       '空翻': 'somersault',
    '曲线向上': 'curve-up', '掉落': 'drop',         '浮动': 'float',
    '基本旋转': 'spin-basic', '螺旋飞入': 'spiral', '玩具风车': 'windmill',
    '字幕式': 'credits',
}

# canonical effect name → (filter_template, default_param, dur_ms)
# {} in template is replaced by param; None means appear-only (no filter)
_EFFECT = {
    'appear':         (None,                    '',            0),
    'blinds':         ('blinds({})',             'horizontal',  500),
    'bounce':         (None,                    '',            0),
    'box':            ('box({})',                'in',          500),
    'center-revolve': (None,                    '',            0),
    'checkerboard':   ('checkboard({})',         'across',      500),
    'circle':         ('circle({})',             'in',          500),
    'compress':       ('stretch(across)',        '',            400),
    'credits':        ('stretch(down)',          '',            600),
    'curve-up':       (None,                    '',            0),
    'diamond':        ('diamond({})',            'in',          500),
    'dissolve':       ('dissolve',               '',            500),
    'drop':           (None,                    '',            0),
    'expand':         ('stretch(across)',        '',            500),
    'fade':           (None,                    '',            400),
    'flip':           (None,                    '',            0),
    'float':          (None,                    '',            0),
    'float-down':     (None,                    '',            0),
    'float-up':       (None,                    '',            0),
    'pinwheel':       ('wheel(1)',               '',            600),
    'plus':           ('plus({})',               'in',          500),
    'random-bars':    ('randombar({})',          'horizontal',  500),
    'rise':           (None,                    '',            0),
    'somersault':     (None,                    '',            0),
    'spin':           ('wheel(1)',               '',            600),
    'spin-basic':     (None,                    '',            0),
    'spiral':         (None,                    '',            0),
    'swivel':         (None,                    '',            0),
    'wedge':          ('wedge',                  '',            500),
    'wheel':          ('wheel({})',              '1',           600),
    'whip':           (None,                    '',            0),
    'windmill':       ('wheel(3)',               '',            600),
    'wipe':           ('wipe({})',               'left',        400),
    'zoom':           (None,                    '',            0),
    'zoom-basic':     (None,                    '',            0),
}

# fly direction → wipe direction (fly enters FROM the opposite side)
_FLY_WIPE = {
    'left': 'right', 'right': 'left', 'up': 'down', 'down': 'up',
    'bottom-left': 'right', 'bottom-right': 'left',
    'top-left': 'right',   'top-right': 'left',
}

# split param → barn filter
_SPLIT_FILTER = {
    'h-out': 'barn(inHorizontal)', 'h-in': 'barn(outHorizontal)',
    'v-out': 'barn(inVertical)',   'v-in': 'barn(outVertical)',
}

# stairs/strips param → camelCase for filter
def _camel(s):
    parts = s.split('-')
    return parts[0] + ''.join(p.capitalize() for p in parts[1:])


def _resolve_anim(spec):
    """
    Parse animation spec (English canonical or legacy Chinese) and return
    (filter_str, dur_ms).  filter_str=None → appear-only, no filter element.
    """
    parts = [p.strip() for p in spec.split(',')]

    # Translate legacy Chinese name → canonical English
    raw_name  = parts[0]
    raw_param = parts[1] if len(parts) > 1 else ''
    name  = _ZH_NAME.get(raw_name, raw_name).lower()
    param = _ZH_PARAM.get(raw_param, raw_param).lower()

    # fly is a special case (wipe proxy)
    if name == 'fly':
        d = _FLY_WIPE.get(param or 'left', 'right')
        return f'wipe({d})', 400

    # split → barn filter
    if name == 'split':
        p = param or 'h-out'
        return _SPLIT_FILTER.get(p, 'barn(inHorizontal)'), 500

    # stairs / strips → parameterised filter
    if name in ('stairs', 'strips'):
        p = _camel(param or 'left-down')
        return f'{name}({p})', 500

    entry = _EFFECT.get(name, (None, '', 0))
    tmpl, default, dur = entry
    if tmpl is None:
        return None, dur
    p = param or default
    return (tmpl.format(p) if '{}' in tmpl else tmpl), dur


def add_entrance_animation(slide, shape, anim_str):
    """
    Append a p:timing element to the slide that auto-plays the entrance
    animation when the slide is displayed.
    Accepts both new English canonical format ("fly,left") and legacy
    Chinese format ("飞入,向左").
    """
    filter_str, dur_ms = _resolve_anim(anim_str)
    preset_id = 1   # Appear preset; filter overrides the visual effect
    spid = str(shape.shape_id)

    i1, i2, i3, i4, i5, i6, i7 = [_next_id() for _ in range(7)]

    P = 'http://schemas.openxmlformats.org/presentationml/2006/main'

    def pe(tag): return '{%s}%s' % (P, tag)

    timing = etree.SubElement(slide._element, pe('timing'))

    # ── tnLst ────────────────────────────────────────────────────────────────
    tnLst = etree.SubElement(timing, pe('tnLst'))
    par0  = etree.SubElement(tnLst,  pe('par'))
    root  = etree.SubElement(par0,   pe('cTn'))
    root.set('id', str(i1)); root.set('dur', 'indefinite')
    root.set('restart', 'whenNotActive'); root.set('nodeType', 'tmRoot')
    child1 = etree.SubElement(root,  pe('childTnLst'))

    seq   = etree.SubElement(child1, pe('seq'))
    seq.set('concurrent', '1'); seq.set('nextAc', 'seek')
    main  = etree.SubElement(seq,    pe('cTn'))
    main.set('id', str(i2)); main.set('dur', 'indefinite')
    main.set('nodeType', 'mainSeq')
    child2 = etree.SubElement(main,  pe('childTnLst'))

    # outer par (auto-start)
    par1  = etree.SubElement(child2, pe('par'))
    ctn3  = etree.SubElement(par1,   pe('cTn'))
    ctn3.set('id', str(i3)); ctn3.set('fill', 'hold')
    sc3   = etree.SubElement(ctn3,   pe('stCondLst'))
    c3    = etree.SubElement(sc3,    pe('cond')); c3.set('delay', '0')
    ch3   = etree.SubElement(ctn3,   pe('childTnLst'))

    par2  = etree.SubElement(ch3,    pe('par'))
    ctn4  = etree.SubElement(par2,   pe('cTn'))
    ctn4.set('id', str(i4)); ctn4.set('fill', 'hold')
    sc4   = etree.SubElement(ctn4,   pe('stCondLst'))
    c4    = etree.SubElement(sc4,    pe('cond')); c4.set('delay', '0')
    ch4   = etree.SubElement(ctn4,   pe('childTnLst'))

    par3  = etree.SubElement(ch4,    pe('par'))
    ctn5  = etree.SubElement(par3,   pe('cTn'))
    ctn5.set('id', str(i5))
    ctn5.set('presetID', str(preset_id)); ctn5.set('presetClass', 'entr')
    ctn5.set('presetSubtype', '0');      ctn5.set('fill', 'hold')
    ctn5.set('grpId', '0');              ctn5.set('nodeType', 'afterEffect')
    sc5  = etree.SubElement(ctn5,    pe('stCondLst'))
    c5   = etree.SubElement(sc5,     pe('cond')); c5.set('delay', '0')
    ch5  = etree.SubElement(ctn5,    pe('childTnLst'))

    # set visibility → visible
    set_e  = etree.SubElement(ch5,   pe('set'))
    cbhvr  = etree.SubElement(set_e, pe('cBhvr'))
    ctn6   = etree.SubElement(cbhvr, pe('cTn'))
    ctn6.set('id', str(i6)); ctn6.set('dur', '1'); ctn6.set('fill', 'hold')
    tgt1   = etree.SubElement(cbhvr, pe('tgtEl'))
    sp1    = etree.SubElement(tgt1,  pe('spTgt')); sp1.set('spid', spid)
    anl    = etree.SubElement(cbhvr, pe('attrNameLst'))
    an     = etree.SubElement(anl,   pe('attrName')); an.text = 'style.visibility'
    to_e   = etree.SubElement(set_e, pe('to'))
    sv     = etree.SubElement(to_e,  pe('strVal')); sv.set('val', 'visible')

    if filter_str:
        eff  = etree.SubElement(ch5, pe('animEffect'))
        eff.set('transition', 'in'); eff.set('filter', filter_str)
        cb2  = etree.SubElement(eff, pe('cBhvr'))
        ctn7 = etree.SubElement(cb2, pe('cTn'))
        ctn7.set('id', str(i7)); ctn7.set('dur', str(dur_ms))
        tgt2 = etree.SubElement(cb2, pe('tgtEl'))
        sp2  = etree.SubElement(tgt2, pe('spTgt')); sp2.set('spid', spid)

    # prevCondLst
    pcl  = etree.SubElement(seq,  pe('prevCondLst'))
    pc   = etree.SubElement(pcl,  pe('cond'))
    pc.set('evt', 'onPrevClick'); pc.set('delay', '0')
    ptn  = etree.SubElement(pc,   pe('tn'))
    pte  = etree.SubElement(ptn,  pe('tnEl'))
    pseq = etree.SubElement(pte,  pe('seq')); pseq.set('nodeType', 'mainSeq')

    # ── bldLst ───────────────────────────────────────────────────────────────
    bldLst = etree.SubElement(timing, pe('bldLst'))
    bldP   = etree.SubElement(bldLst, pe('bldP'))
    bldP.set('spid', spid); bldP.set('grpId', '0')
    bldP.set('uiExpand', '1'); bldP.set('build', 'byChar')


# ---------------------------------------------------------------------------
# Core layout logic
# ---------------------------------------------------------------------------
def rotation_delta(enter_rot_str):
    if enter_rot_str == 'left':  return -90
    if enter_rot_str == 'right': return +90
    return 0


def pptx_rot(deg):
    return int(deg % 360)


def generate_slide(prs, slide_entries, current_phase_idx, last_trans_dir,
                   bg_path, blur_radius, default_font, default_color,
                   newest_anim):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    if bg_path and os.path.isfile(bg_path):
        add_background(slide, bg_path, blur_radius)

    current = [e for e in slide_entries if e['phase_idx'] == current_phase_idx]
    old     = [e for e in slide_entries if e['phase_idx'] != current_phase_idx]

    # ── Active phase: newest at BOTTOM_CM, previous lines stacked above ──────
    y_bottom = BOTTOM_CM
    newest_shape = None
    for i, entry in enumerate(reversed(current)):   # newest first
        h      = line_height_cm(entry['font_size_pt'])
        top_cm = y_bottom - h
        shape  = add_textbox(slide, entry['text'], default_font, default_color,
                             entry['font_size_pt'],
                             LEFT_CM, top_cm, BOX_W_CM, h)
        if i == 0:
            newest_shape = shape   # track the newest for animation
        y_bottom = top_cm - GAP_CM

    if newest_shape and newest_anim:
        add_entrance_animation(slide, newest_shape, newest_anim)

    # ── Old phases: packed outward from the slide edge ────────────────────────
    if not old:
        return slide

    phase_groups = {}
    for e in old:
        phase_groups.setdefault(e['phase_idx'], []).append(e)

    side        = last_trans_dir
    cy          = SLIDE_H_CM / 2.0
    inner_limit = 0.0 if side == 'left' else SLIDE_W_CM

    for pi in sorted(phase_groups.keys(), reverse=True):
        group    = phase_groups[pi]
        rot      = group[0]['phase_rotation']
        rot_norm = rot % 360

        if rot_norm in (90, 270):
            for entry in reversed(group):   # newest first → closest to slide
                col_h    = line_height_cm(entry['font_size_pt'])
                rot_pptx = pptx_rot(rot)
                if side == 'left':
                    center_x    = inner_limit - col_h / 2.0
                    inner_limit -= col_h
                else:
                    center_x    = inner_limit + col_h / 2.0
                    inner_limit += col_h
                add_textbox(slide, entry['text'], default_font, default_color,
                            entry['font_size_pt'],
                            center_x - BOX_W_CM / 2.0, cy - col_h / 2.0,
                            BOX_W_CM, col_h, rotation_deg=rot_pptx)
        else:
            y_b = BOTTOM_CM
            for entry in reversed(group):
                h      = line_height_cm(entry['font_size_pt'])
                top_cm = y_b - h
                left_cm = (inner_limit - BOX_W_CM) if side == 'left' \
                          else inner_limit
                add_textbox(slide, entry['text'], default_font, default_color,
                            entry['font_size_pt'],
                            left_cm, top_cm, BOX_W_CM, h)
                y_b = top_cm - GAP_CM
            inner_limit += (-BOX_W_CM if side == 'left' else BOX_W_CM)

    return slide


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------
def generate_pptx(json_path, output_path):
    base_dir = os.path.dirname(json_path)

    with open(json_path, encoding='utf-8') as f:
        config = json.load(f)

    defaults      = config.get('defaults', {})
    default_font  = defaults.get('font', '黑体')
    default_color = defaults.get('color', '#000000')
    bg_file       = defaults.get('background', {}).get('file', '')
    bg_blur       = defaults.get('background', {}).get('blur', 0)
    bg_path       = os.path.join(base_dir, bg_file) if bg_file else ''

    phases = config['phases']

    prs = Presentation()
    prs.slide_width  = cm(SLIDE_W_CM)
    prs.slide_height = cm(SLIDE_H_CM)

    all_entries    = []
    last_trans_dir = 'none'

    for phase_idx, phase in enumerate(phases):
        enter_rot_str = phase.get('enter_rotation', 'none')
        delta         = rotation_delta(enter_rot_str)

        if delta != 0:
            for e in all_entries:
                e['phase_rotation'] += delta
            last_trans_dir = enter_rot_str

        # Discard oldest old phase when there are more than MAX_OLD_PHASES
        if phase_idx > 0:
            old_phases = sorted(set(e['phase_idx'] for e in all_entries))
            while len(old_phases) > MAX_OLD_PHASES:
                oldest = old_phases.pop(0)
                all_entries = [e for e in all_entries
                               if e['phase_idx'] != oldest]

        for line_idx, srt_entry in enumerate(phase['srt']):
            _start, _end, text, anim_type = srt_entry

            entry = {
                'text':           text,
                'font_size_pt':   BASE_PT,
                'phase_idx':      phase_idx,
                'phase_rotation': 0,
                'anim_type':      anim_type,
            }
            all_entries.append(entry)

            # Compute font sizes for every line in the current phase.
            # Newest line: natural size from text length.
            # Older lines: scale by SCALE_BACK per step back (camera zoom).
            cur = [e for e in all_entries if e['phase_idx'] == phase_idx]
            n   = len(cur)
            n0  = natural_font_size(cur[-1]['text'])   # newest natural size
            for k, e in enumerate(cur):
                steps_back = (n - 1) - k
                e['font_size_pt'] = n0 * (SCALE_BACK ** steps_back)

            generate_slide(prs, all_entries, phase_idx, last_trans_dir,
                           bg_path, bg_blur, default_font, default_color,
                           anim_type)

    prs.save(output_path)
    print(f"Saved: {output_path}  ({len(prs.slides)} slides)")


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print("Usage: generate_pptx.py <input.json> <output.pptx>")
        sys.exit(1)
    generate_pptx(sys.argv[1], sys.argv[2])
